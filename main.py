from fastapi import FastAPI, HTTPException
from fastapi.responses import JSONResponse, FileResponse
from pydantic import BaseModel, HttpUrl
from typing import List, Optional, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os, uuid, io, requests

# Image conversion
from PIL import Image

# Optional SVG support
try:
    import cairosvg
    HAS_CAIROSVG = True
except Exception:
    HAS_CAIROSVG = False

app = FastAPI()

# -------------------------
# Themes available on disk
# -------------------------
# Put your .potx files in ./themes and list them here (key -> filename)
THEMES: Dict[str, str] = {
    # examples — add your own .potx files to the repo
    "corporate_blue": "themes/corporate_blue.potx",
    "dark_minimal": "themes/dark_minimal.potx",
    "clean": "themes/clean.potx",
}

# -------------------------
# Health
# -------------------------
@app.get("/")
def health():
    return {
        "status": "ok",
        "endpoints": ["/docs", "/pptx/create"],
        "themes": list(THEMES.keys()),
    }

# -------------------------
# Models
# -------------------------
class ImageItem(BaseModel):
    url: HttpUrl
    width_inch: Optional[float] = None
    height_inch: Optional[float] = None
    caption: Optional[str] = None

class SlideItem(BaseModel):
    heading: str
    bullets: Optional[List[str]] = []
    images: Optional[List[ImageItem]] = []

class CreatePptxInput(BaseModel):
    title: str
    subtitle: Optional[str] = None
    slides: List[SlideItem]

    # NEW — theming & branding
    theme: Optional[str] = None           # e.g., "corporate_blue"
    theme_url: Optional[HttpUrl] = None   # load a .potx by URL
    primary_color: Optional[str] = None   # hex "#1E88E5" or "1E88E5"
    dark_mode: Optional[bool] = False
    logo_url: Optional[HttpUrl] = None    # top-right logo on all slides

    footer: Optional[str] = None

# -------------------------
# Storage
# -------------------------
FILES_DIR = "generated"
os.makedirs(FILES_DIR, exist_ok=True)

# -------------------------
# Helpers
# -------------------------
def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.strip().lstrip("#")
    if len(h) == 3:
        h = "".join([c*2 for c in h])
    if len(h) != 6:
        raise ValueError("primary_color must be a 3/6-digit hex like #1E88E5")
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    return RGBColor(r, g, b)

def fetch_bytes(url: str) -> bytes:
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; PPTX-Generator/1.3)",
        "Accept": "*/*",
    }
    r = requests.get(url, headers=headers, timeout=25, allow_redirects=True)
    r.raise_for_status()
    return r.content

def fetch_image_bytes(url: str) -> io.BytesIO:
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; PPTX-Generator/1.3)",
        "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
    }
    r = requests.get(url, headers=headers, timeout=25, allow_redirects=True)
    r.raise_for_status()
    ct = (r.headers.get("Content-Type") or "").lower()
    url_l = url.lower()

    # SVG → PNG (optional)
    if ("image/svg" in ct) or url_l.endswith(".svg"):
        if not HAS_CAIROSVG:
            raise ValueError("SVG found but cairosvg not installed on server.")
        return io.BytesIO(cairosvg.svg2png(bytestring=r.content))

    # WebP → PNG
    if ("image/webp" in ct) or url_l.endswith(".webp"):
        img = Image.open(io.BytesIO(r.content)).convert("RGBA")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        return buf

    # PNG/JPEG/GIF
    if any(t in ct for t in ("image/png", "image/jpeg", "image/jpg", "image/gif")):
        return io.BytesIO(r.content)
    if url_l.endswith((".png", ".jpg", ".jpeg", ".gif")):
        return io.BytesIO(r.content)

    raise ValueError(f"Unsupported image type: {ct or 'unknown'}")

def _apply_background(slide, rgb: Optional[RGBColor], dark: bool):
    try:
        fill = slide.background.fill
        fill.solid()
        if rgb:
            fill.fore_color.rgb = rgb
        else:
            # subtle defaults if not provided
            fill.fore_color.rgb = RGBColor(18, 18, 18) if dark else RGBColor(255, 255, 255)
    except Exception:
        pass

def _style_title(slide, rgb: Optional[RGBColor], dark: bool):
    try:
        title = slide.shapes.title
        if not title:
            return
        run = title.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(40)
        if rgb:
            run.font.color.rgb = rgb
        elif dark:
            run.font.color.rgb = RGBColor(230, 230, 230)
    except Exception:
        pass

def _add_logo(slide, prs: Presentation, logo_stream: Optional[io.BytesIO]):
    if not logo_stream:
        return
    try:
        # top-right; 1.1" width, maintain aspect
        pic = slide.shapes.add_picture(logo_stream, prs.slide_width - Inches(1.6), Inches(0.2), width=Inches(1.1))
        return pic
    except Exception:
        pass

# -------------------------
# PPTX Builder
# -------------------------
def build_pptx(payload: CreatePptxInput, output_path: str):
    # --- Choose a template/theme ---
    prs: Presentation
    if payload.theme_url:
        prs = Presentation(io.BytesIO(fetch_bytes(str(payload.theme_url))))
    elif payload.theme and payload.theme in THEMES and os.path.exists(THEMES[payload.theme]):
        prs = Presentation(THEMES[payload.theme])
    else:
        prs = Presentation()

    # Precompute style params
    rgb = _hex_to_rgb(payload.primary_color) if payload.primary_color else None
    logo_stream = io.BytesIO(fetch_bytes(str(payload.logo_url))) if payload.logo_url else None

    # Layouts
    title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]

    # --- Title slide ---
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = payload.title

    # background & title styling
    _apply_background(slide, rgb=None if not payload.dark_mode else RGBColor(18, 18, 18), dark=bool(payload.dark_mode))
    _style_title(slide, rgb, bool(payload.dark_mode))
    _add_logo(slide, prs, logo_stream)

    if payload.subtitle:
        try:
            slide.placeholders[1].text = payload.subtitle
            # subtitle styling
            sub_tf = slide.placeholders[1].text_frame
            sub_tf.paragraphs[0].runs[0].font.size = Pt(20)
            if payload.dark_mode:
                sub_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(210, 210, 210)
        except Exception:
            pass

    # --- Content slides ---
    for s in payload.slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = s.heading[:255]
        _apply_background(slide, rgb=None if not payload.dark_mode else RGBColor(18, 18, 18), dark=bool(payload.dark_mode))
        _style_title(slide, rgb, bool(payload.dark_mode))
        _add_logo(slide, prs, logo_stream)

        # Bullets
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        if s.bullets:
            body.text = s.bullets[0][:1000]
            for b in s.bullets[1:]:
                p = body.add_paragraph()
                p.text = b[:1000]
                p.level = 0
            # bullet color for dark mode
            if payload.dark_mode:
                try:
                    for para in body.paragraphs:
                        for run in para.runs:
                            run.font.color.rgb = RGBColor(230, 230, 230)
                except Exception:
                    pass

        # Images
        if s.images:
            top = Inches(2.8)
            max_width = Inches(6.5)
            for img in s.images:
                try:
                    stream = fetch_image_bytes(str(img.url))
                    if img.width_inch and img.height_inch:
                        pic = slide.shapes.add_picture(stream, Inches(0.5), top,
                                                       width=Inches(img.width_inch),
                                                       height=Inches(img.height_inch))
                    elif img.width_inch:
                        pic = slide.shapes.add_picture(stream, Inches(0.5), top, width=Inches(img.width_inch))
                    elif img.height_inch:
                        pic = slide.shapes.add_picture(stream, Inches(0.5), top, height=Inches(img.height_inch))
                    else:
                        pic = slide.shapes.add_picture(stream, Inches(1), top, width=max_width)

                    pic.left = int((prs.slide_width - pic.width) / 2)

                    if img.caption:
                        cap = slide.shapes.add_textbox(pic.left, pic.top + pic.height + Inches(0.08),
                                                       pic.width, Inches(0.36))
                        cap_tf = cap.text_frame
                        cap_tf.text = img.caption[:200]
                        try:
                            cap_tf.paragraphs[0].runs[0].font.size = Pt(12)
                            if payload.dark_mode:
                                cap_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(220, 220, 220)
                        except Exception:
                            pass

                    top = pic.top + pic.height + Inches(0.24)

                except Exception as e:
                    p = body.add_paragraph()
                    p.text = f"[Image failed: {img.url} — {str(e)}]"
                    p.level = 0

    # Footer
    if payload.footer:
        for sld in prs.slides:
            tx = sld.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.3))
            run = tx.text_frame.paragraphs[0].add_run()
            run.text = payload.footer[:120]
            run.font.size = Pt(10)
            if payload.dark_mode:
                run.font.color.rgb = RGBColor(200, 200, 200)

    prs.save(output_path)

# -------------------------
# API Routes
# -------------------------
@app.post("/pptx/create")
async def create_pptx(payload: CreatePptxInput):
    file_id = uuid.uuid4().hex
    filename = f"{file_id}.pptx"
    path = os.path.join(FILES_DIR, filename)

    build_pptx(payload, path)

    base_url = os.getenv("PUBLIC_BASE_URL", "http://localhost:8000")
    return JSONResponse({"download_url": f"{base_url}/files/{filename}", "file_name": filename})

@app.get("/files/{filename}")
async def serve_file(filename: str):
    path = os.path.join(FILES_DIR, filename)
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(
        path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )
